from trello import TrelloClient
from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import *

boardID = "bn3F6E4m" # whatever board that is, eg. https://trello.com/b/uDtSJilF/iab207
#init
client = TrelloClient(
    api_key='c75ebe201687acd0b1d4c55ec98aa42f',
    api_secret='02855a09c9756eb487c4e83b76ff5e3090ad111860b3cdbd0c5476d0b1fcd7e6',
    token='b03a2e5dd545cc1fbeebc65e520aafaa0faf6a37eee4df6d3a1e21e7da16475a'
)

# get board
board = client.get_board(boardID)

# tap that list
list_sections = board.list_lists()

# Generate Presentation
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)

left = top = width = height = Inches(1)
background = slide.shapes.add_picture(r'static\background.png', 0, 0)

title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "QUT MOTORSPORT"
subtitle.text = input("What would you like the week/subtitle to be?\n")

txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame

today = date.today()
comp = date(2019,12,5)
diff = comp - today
tf.text = str(diff.days) + " Days til Comp"

# Slides
for num, lead in enumerate(list_sections):

    dt = datetime.now(timezone.utc)
    is_content = False

    for card in lead.list_cards():
        try:

            due_diff = dt - card.due_date
            due_days = due_diff.days * -1
            
            create_diff = dt - card.created_date
            
            if (due_days <= 14 or create_diff.days <= 7):
                is_content = True

        except:

            create_diff = dt - card.created_date
            
            if (create_diff.days <= 7):
                is_content = True

    if (is_content):
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = lead.name
        
        for card in lead.list_cards():
            
            print(card.name + " " + str(card.created_date) + " " + str(card.due_date))
            
            try:

                due_diff = dt - card.due_date
                due_days = due_diff.days * -1
                
                create_diff = dt - card.created_date
                
                if (due_days <= 14 or create_diff.days <= 7):
                    print("Show")

                    tf = body_shape.text_frame
                    t = tf.add_paragraph()
                    t.text = card.name

                    p = tf.add_paragraph()
                    p.text = card.desc
                    p.level = 1

            except:

                create_diff = dt - card.created_date
                
                if (create_diff.days <= 7):
                    print("Show")

                    tf = body_shape.text_frame
                    t = tf.add_paragraph()
                    t.text = card.name

                    p = tf.add_paragraph()
                    p.text = card.desc
                    p.level = 1

            # print(dir(card))

        # Debugging

        # for shape in slide.placeholders:
        #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
        
        # # print(dir(lead))

prs.save('test.pptx')
